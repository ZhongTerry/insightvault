"""
多模态解析模块 (Multi-modal Parser)
支持 PDF、Word、PPT、图片、音视频等多种格式的内容提取
"""

import os
import io
import base64
import logging
import mimetypes
from pathlib import Path
from typing import Dict, Optional, List, Any
from dataclasses import dataclass

# PDF 解析
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

# Word 文档解析
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PowerPoint 解析
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 图片处理
try:
    from PIL import Image
    import cv2
    import numpy as np
    HAS_IMAGE = True
except ImportError:
    HAS_IMAGE = False

# OCR
try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

# 音频转写
try:
    import whisper
    HAS_WHISPER = True
except ImportError:
    HAS_WHISPER = False

logger = logging.getLogger("InsightVault.MultiModal")


@dataclass
class ParseResult:
    """解析结果数据类"""
    success: bool
    content: str  # 提取的文本内容
    metadata: Dict[str, Any]  # 元数据（作者、创建时间等）
    images: List[Dict[str, str]] = None  # 提取的图片（base64）
    error: Optional[str] = None


class MultiModalParser:
    """多模态文件解析器"""
    
    def __init__(self, ocr_language: str = 'chi_sim+eng'):
        """
        初始化解析器
        
        Args:
            ocr_language: Tesseract OCR 语言包（默认中英文）
        """
        self.ocr_language = ocr_language
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查依赖库是否安装"""
        status = {
            "PDF (PyMuPDF)": HAS_PYMUPDF,
            "Word (python-docx)": HAS_DOCX,
            "PowerPoint (python-pptx)": HAS_PPTX,
            "Image Processing (PIL/OpenCV)": HAS_IMAGE,
            "OCR (Tesseract)": HAS_TESSERACT,
            "Audio (Whisper)": HAS_WHISPER
        }
        
        for feature, available in status.items():
            if available:
                logger.info(f"✓ {feature} 可用")
            else:
                logger.warning(f"✗ {feature} 不可用（未安装依赖）")
    
    def parse_file(self, file_path: str, mime_type: Optional[str] = None) -> ParseResult:
        """
        自动识别并解析文件
        
        Args:
            file_path: 文件路径
            mime_type: MIME 类型（可选，自动检测）
        
        Returns:
            ParseResult: 解析结果
        """
        if not os.path.exists(file_path):
            return ParseResult(success=False, content="", metadata={}, error="文件不存在")
        
        if mime_type is None:
            mime_type, _ = mimetypes.guess_type(file_path)
        
        logger.info(f"开始解析文件: {file_path} (类型: {mime_type})")
        
        # 根据 MIME 类型分发到不同的解析器
        if mime_type == "application/pdf":
            return self.parse_pdf(file_path)
        elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
            return self.parse_word(file_path)
        elif mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
            return self.parse_ppt(file_path)
        elif mime_type and mime_type.startswith("image/"):
            return self.parse_image(file_path)
        elif mime_type and mime_type.startswith("audio/"):
            return self.parse_audio(file_path)
        elif mime_type and mime_type.startswith("video/"):
            return self.parse_video(file_path)
        else:
            # 尝试作为纯文本读取
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return ParseResult(success=True, content=content, metadata={"type": "text"})
            except Exception as e:
                return ParseResult(success=False, content="", metadata={}, error=f"不支持的文件类型: {mime_type}")
    
    def parse_pdf(self, file_path: str) -> ParseResult:
        """解析 PDF 文件"""
        if not HAS_PYMUPDF:
            return ParseResult(success=False, content="", metadata={}, error="PyMuPDF 未安装")
        
        try:
            doc = fitz.open(file_path)
            content_parts = []
            images = []
            metadata = {
                "page_count": doc.page_count,
                "author": doc.metadata.get("author", ""),
                "title": doc.metadata.get("title", ""),
                "subject": doc.metadata.get("subject", ""),
                "creator": doc.metadata.get("creator", ""),
                "producer": doc.metadata.get("producer", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
            }
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # 提取文本
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")
                
                # 提取图片（如果需要）
                for img_idx, img in enumerate(page.get_images()):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        image_b64 = base64.b64encode(image_bytes).decode()
                        images.append({
                            "page": page_num + 1,
                            "index": img_idx,
                            "format": image_ext,
                            "data": image_b64
                        })
                    except Exception as e:
                        logger.warning(f"提取图片失败 (页 {page_num + 1}, 图 {img_idx}): {e}")
            
            doc.close()
            
            full_content = "\n\n".join(content_parts)
            return ParseResult(
                success=True,
                content=full_content,
                metadata=metadata,
                images=images
            )
        
        except Exception as e:
            logger.error(f"PDF 解析失败: {e}")
            return ParseResult(success=False, content="", metadata={}, error=str(e))
    
    def parse_word(self, file_path: str) -> ParseResult:
        """解析 Word 文档"""
        if not HAS_DOCX:
            return ParseResult(success=False, content="", metadata={}, error="python-docx 未安装")
        
        try:
            doc = Document(file_path)
            content_parts = []
            
            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            
            # 提取表格
            for table_idx, table in enumerate(doc.tables):
                table_text = f"\n--- 表格 {table_idx + 1} ---\n"
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_text += row_text + "\n"
                content_parts.append(table_text)
            
            # 元数据
            core_props = doc.core_properties
            metadata = {
                "author": core_props.author or "",
                "title": core_props.title or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
            }
            
            full_content = "\n\n".join(content_parts)
            return ParseResult(success=True, content=full_content, metadata=metadata)
        
        except Exception as e:
            logger.error(f"Word 解析失败: {e}")
            return ParseResult(success=False, content="", metadata={}, error=str(e))
    
    def parse_ppt(self, file_path: str) -> ParseResult:
        """解析 PowerPoint 文档"""
        if not HAS_PPTX:
            return ParseResult(success=False, content="", metadata={}, error="python-pptx 未安装")
        
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = f"--- 幻灯片 {slide_idx + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                content_parts.append(slide_text)
            
            # 元数据
            core_props = prs.core_properties
            metadata = {
                "author": core_props.author or "",
                "title": core_props.title or "",
                "subject": core_props.subject or "",
                "slide_count": len(prs.slides),
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
            }
            
            full_content = "\n\n".join(content_parts)
            return ParseResult(success=True, content=full_content, metadata=metadata)
        
        except Exception as e:
            logger.error(f"PPT 解析失败: {e}")
            return ParseResult(success=False, content="", metadata={}, error=str(e))
    
    def parse_image(self, file_path: str, use_local_ocr: bool = True) -> ParseResult:
        """
        解析图片并提取文字
        
        Args:
            file_path: 图片路径
            use_local_ocr: 是否使用本地 Tesseract OCR（否则返回 base64 供 AI OCR）
        """
        if not HAS_IMAGE:
            return ParseResult(success=False, content="", metadata={}, error="PIL/OpenCV 未安装")
        
        try:
            # 读取图片
            img = Image.open(file_path)
            metadata = {
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "mode": img.mode
            }
            
            # 转换为 base64
            buffered = io.BytesIO()
            img.save(buffered, format=img.format or "PNG")
            img_b64 = base64.b64encode(buffered.getvalue()).decode()
            
            content = ""
            
            # 本地 OCR
            if use_local_ocr and HAS_TESSERACT:
                try:
                    # 预处理图片提高 OCR 准确率
                    img_cv = cv2.imread(file_path)
                    gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
                    # 自适应阈值二值化
                    processed = cv2.adaptiveThreshold(
                        gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
                    )
                    
                    # OCR
                    text = pytesseract.image_to_string(processed, lang=self.ocr_language)
                    content = text.strip()
                    metadata["ocr_method"] = "tesseract"
                except Exception as e:
                    logger.warning(f"本地 OCR 失败: {e}")
                    metadata["ocr_method"] = "none"
            else:
                metadata["ocr_method"] = "none"
            
            return ParseResult(
                success=True,
                content=content,
                metadata=metadata,
                images=[{"data": img_b64, "format": img.format}]
            )
        
        except Exception as e:
            logger.error(f"图片解析失败: {e}")
            return ParseResult(success=False, content="", metadata={}, error=str(e))
    
    def parse_audio(self, file_path: str) -> ParseResult:
        """解析音频文件（使用 Whisper 转写）"""
        if not HAS_WHISPER:
            return ParseResult(
                success=False, 
                content="", 
                metadata={}, 
                error="Whisper 未安装（大型依赖，需手动安装: pip install openai-whisper）"
            )
        
        try:
            # 加载 Whisper 模型（首次会下载）
            model = whisper.load_model("base")  # 可选: tiny, base, small, medium, large
            
            # 转写音频
            result = model.transcribe(file_path, language="zh")  # 可自动检测语言
            
            metadata = {
                "duration": result.get("duration", 0),
                "language": result.get("language", "unknown")
            }
            
            # 包含时间戳的详细转写
            segments = []
            for seg in result.get("segments", []):
                segments.append(f"[{seg['start']:.2f}s - {seg['end']:.2f}s] {seg['text']}")
            
            content = result["text"]  # 完整文本
            detailed_content = "\n".join(segments)  # 带时间戳的版本
            
            metadata["detailed_transcript"] = detailed_content
            
            return ParseResult(success=True, content=content, metadata=metadata)
        
        except Exception as e:
            logger.error(f"音频解析失败: {e}")
            return ParseResult(success=False, content="", metadata={}, error=str(e))
    
    def parse_video(self, file_path: str) -> ParseResult:
        """解析视频文件（提取音轨并转写）"""
        # 提取音轨需要 ffmpeg，简化版本暂不实现
        return ParseResult(
            success=False,
            content="",
            metadata={},
            error="视频解析需要 ffmpeg 支持，建议先提取音轨后使用音频解析"
        )


# 工厂函数
def create_parser(ocr_language: str = 'chi_sim+eng') -> MultiModalParser:
    """创建多模态解析器实例"""
    return MultiModalParser(ocr_language=ocr_language)


# 便捷函数
def parse_file(file_path: str, mime_type: Optional[str] = None) -> ParseResult:
    """快速解析文件"""
    parser = create_parser()
    return parser.parse_file(file_path, mime_type)
